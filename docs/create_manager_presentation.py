from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slides = [
    ("Playwright API Automation Framework", ["Using TypeScript and GitHub Copilot", "Engineering Management Demo", "Lightweight, reusable, and CI-ready"]),
    ("Project Objective", ["Learn modern API automation design", "Validate CRUD operations reliably", "Deliver a reusable framework for future scale"]),
    ("Technology Stack", ["Playwright", "TypeScript", "Node.js", "GitHub Actions", "JSONPlaceholder API"]),
    ("Framework Architecture", ["Test Layer", "Service Layer", "API Client Layer", "External API"]),
    ("Framework Implementation", ["ApiClient", "Service Layer", "Test Layer", "Configuration"]),
    ("CRUD Automation", ["Create", "Read", "Update", "Delete", "Status and payload validation"]),
    ("Reporting and CI/CD", ["HTML report", "JSON report", "JUnit report", "GitHub Actions workflow"]),
    ("Production Readiness", ["Environment management", "Error handling", "Logging", "CI improvements"]),
    ("GitHub Copilot Experience", ["Prompt-driven development", "Faster implementation", "Documentation generation", "Iterative refinement"]),
    ("Future Roadmap and Q&A", ["Authentication management", "Schema validation", "Secrets management", "Enterprise integration"]),
]

for title, bullets in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Calibri'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 54, 93)

    body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(11.5), Inches(4.8))
    tfb = body_box.text_frame
    tfb.word_wrap = True
    tfb.clear()
    for index, bullet in enumerate(bullets):
        p = tfb.paragraphs[0] if index == 0 else tfb.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = 'Calibri'
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT

prs.save('Manager_Demo_Presentation.pptx')
print('Saved Manager_Demo_Presentation.pptx')
